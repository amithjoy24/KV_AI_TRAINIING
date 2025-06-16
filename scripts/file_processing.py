import os
from io import BytesIO
from PIL import Image
import fitz  # PyMuPDF
import docx
from pptx import Presentation

from scripts.ocr_utils import ocr_image_from_bytes  # use OCR utils

def extract_text_images_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    images = []

    for page_num in range(len(doc)):
        page = doc[page_num]
        text += page.get_text()

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]

            ocr_text = ocr_image_from_bytes(img_bytes)

            images.append({
                "index": f"{page_num}_{img_index}",
                "ocr_text": ocr_text,
                "image_bytes": img_bytes,
            })

    return text, images


def extract_text_images_from_docx(file_path):
    doc = docx.Document(file_path)
    text = "\n".join(p.text for p in doc.paragraphs)
    images = []

    for rel in doc.part._rels:
        rel_obj = doc.part._rels[rel]
        # Check if this rel is an image
        if "image" in rel_obj.target_ref:
            image_data = rel_obj.target_part.blob

            ocr_text = ocr_image_from_bytes(image_data)

            images.append({
                "index": rel,
                "ocr_text": ocr_text,
                "image_bytes": image_data,
            })

    return text, images


def extract_text_images_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    images = []

    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"
            # shape_type 13 means picture
            if shape.shape_type == 13:
                img = shape.image
                image_bytes = img.blob

                ocr_text = ocr_image_from_bytes(image_bytes)

                images.append({
                    "index": f"{slide_num}",
                    "ocr_text": ocr_text,
                    "image_bytes": image_bytes,
                })

    return text, images


def extract_from_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".pdf":
        return extract_text_images_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_images_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_images_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")
